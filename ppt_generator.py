#!/usr/bin/env python3
"""
PowerPoint Presentation Generator

A comprehensive tool for creating PowerPoint presentations programmatically.
Perfect for DevOps, technical documentation, and automated report generation.

Author: AI Assistant
Date: 2025
"""

import os
import sys
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.slide import Slide
    from pptx.util import Inches, Pt
    from pptx.shapes.picture import Picture
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.dml import MSO_THEME_COLOR
    from PIL import Image
except ImportError as e:
    print(f"Error: Required libraries not found. Please install them:")
    print("pip install python-pptx Pillow")
    sys.exit(1)


class PPTGenerator:
    """
    A comprehensive PowerPoint presentation generator with support for:
    - Title slides
    - Content slides with bullet points
    - Image slides
    - Chart slides (bar, line, pie)
    - Two-column layouts
    - Custom formatting and themes
    """
    
    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize the PPT generator.
        
        Args:
            template_path: Optional path to a PowerPoint template file
        """
        if template_path and os.path.exists(template_path):
            self.presentation = Presentation(template_path)
        else:
            self.presentation = Presentation()
        
        # Default styling
        self.title_font_size = Pt(44)
        self.subtitle_font_size = Pt(28)
        self.content_font_size = Pt(18)
        self.primary_color = RGBColor(31, 73, 125)  # Dark blue
        self.accent_color = RGBColor(79, 129, 189)  # Light blue
        
    def add_title_slide(self, title: str, subtitle: str = "", author: str = "") -> Slide:
        """
        Add a title slide to the presentation.
        
        Args:
            title: Main title text
            subtitle: Subtitle text
            author: Author name
            
        Returns:
            The created slide object
        """
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = self.title_font_size
        title_paragraph.font.color.rgb = self.primary_color
        title_paragraph.font.bold = True
        
        # Set subtitle
        if len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            full_subtitle = subtitle
            if author:
                full_subtitle += f"\n\nPresented by: {author}"
            if full_subtitle:
                subtitle_shape.text = full_subtitle
                subtitle_paragraph = subtitle_shape.text_frame.paragraphs[0]
                subtitle_paragraph.font.size = self.subtitle_font_size
                subtitle_paragraph.font.color.rgb = self.accent_color
        
        return slide
    
    def add_content_slide(self, title: str, content: List[str], 
                         layout_type: str = "bullet") -> Slide:
        """
        Add a content slide with bullet points or numbered list.
        
        Args:
            title: Slide title
            content: List of content items
            layout_type: "bullet" or "numbered"
            
        Returns:
            The created slide object
        """
        slide_layout = self.presentation.slide_layouts[1]  # Title and Content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(36)
        title_paragraph.font.color.rgb = self.primary_color
        title_paragraph.font.bold = True
        
        # Add content
        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for i, item in enumerate(content):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = item
                p.font.size = self.content_font_size
                p.font.color.rgb = RGBColor(64, 64, 64)
                p.level = 0
                
                if layout_type == "numbered":
                    p.text = f"{i+1}. {item}"
        
        return slide
    
    def add_two_column_slide(self, title: str, left_content: List[str], 
                           right_content: List[str]) -> Slide:
        """
        Add a slide with two-column layout.
        
        Args:
            title: Slide title
            left_content: Content for left column
            right_content: Content for right column
            
        Returns:
            The created slide object
        """
        slide_layout = self.presentation.slide_layouts[3]  # Two Content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(36)
        title_paragraph.font.color.rgb = self.primary_color
        title_paragraph.font.bold = True
        
        # Add left column content
        if len(slide.placeholders) > 1:
            left_shape = slide.placeholders[1]
            left_frame = left_shape.text_frame
            left_frame.clear()
            
            for i, item in enumerate(left_content):
                p = left_frame.paragraphs[0] if i == 0 else left_frame.add_paragraph()
                p.text = item
                p.font.size = self.content_font_size
                p.font.color.rgb = RGBColor(64, 64, 64)
        
        # Add right column content
        if len(slide.placeholders) > 2:
            right_shape = slide.placeholders[2]
            right_frame = right_shape.text_frame
            right_frame.clear()
            
            for i, item in enumerate(right_content):
                p = right_frame.paragraphs[0] if i == 0 else right_frame.add_paragraph()
                p.text = item
                p.font.size = self.content_font_size
                p.font.color.rgb = RGBColor(64, 64, 64)
        
        return slide
    
    def add_image_slide(self, title: str, image_path: str, 
                       caption: str = "") -> Slide:
        """
        Add a slide with an image.
        
        Args:
            title: Slide title
            image_path: Path to the image file
            caption: Optional image caption
            
        Returns:
            The created slide object
        """
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add title
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(36)
        title_paragraph.font.color.rgb = self.primary_color
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add image if it exists
        if os.path.exists(image_path):
            left = Inches(2)
            top = Inches(2)
            slide.shapes.add_picture(image_path, left, top, width=Inches(6))
            
            # Add caption if provided
            if caption:
                caption_left = Inches(1)
                caption_top = Inches(6.5)
                caption_width = Inches(8)
                caption_height = Inches(0.5)
                
                caption_box = slide.shapes.add_textbox(caption_left, caption_top, 
                                                     caption_width, caption_height)
                caption_frame = caption_box.text_frame
                caption_frame.text = caption
                caption_paragraph = caption_frame.paragraphs[0]
                caption_paragraph.font.size = Pt(14)
                caption_paragraph.font.color.rgb = RGBColor(96, 96, 96)
                caption_paragraph.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_chart_slide(self, title: str, chart_data: Dict[str, Any], 
                       chart_type: str = "column") -> Slide:
        """
        Add a slide with a chart.
        
        Args:
            title: Slide title
            chart_data: Dictionary containing chart data
                       Format: {"categories": [...], "series": [{"name": "...", "values": [...]}]}
            chart_type: Type of chart ("column", "line", "pie")
            
        Returns:
            The created slide object
        """
        slide_layout = self.presentation.slide_layouts[5]  # Title and Content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(36)
        title_paragraph.font.color.rgb = self.primary_color
        title_paragraph.font.bold = True
        
        # Prepare chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = chart_data.get("categories", [])
        
        for series in chart_data.get("series", []):
            chart_data_obj.add_series(series["name"], series["values"])
        
        # Add chart
        left = Inches(2)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4.5)
        
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE
        }
        
        chart_obj = slide.shapes.add_chart(
            chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
            left, top, width, height, chart_data_obj
        )
        
        return slide
    
    def add_section_slide(self, section_title: str, 
                         background_color: Optional[Tuple[int, int, int]] = None) -> Slide:
        """
        Add a section divider slide.
        
        Args:
            section_title: Title of the section
            background_color: Optional RGB tuple for background color
            
        Returns:
            The created slide object
        """
        slide_layout = self.presentation.slide_layouts[2]  # Section header layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set section title
        title_shape = slide.shapes.title
        title_shape.text = section_title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(54)
        title_paragraph.font.color.rgb = RGBColor(255, 255, 255)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Set background color if provided
        if background_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*background_color)
        
        return slide
    
    def save(self, filename: str) -> str:
        """
        Save the presentation to a file.
        
        Args:
            filename: Output filename (with .pptx extension)
            
        Returns:
            Full path to the saved file
        """
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        
        full_path = os.path.abspath(filename)
        self.presentation.save(full_path)
        return full_path
    
    def get_slide_count(self) -> int:
        """Get the number of slides in the presentation."""
        return len(self.presentation.slides)
    
    def set_theme_colors(self, primary_color: Tuple[int, int, int], 
                        accent_color: Tuple[int, int, int]):
        """
        Set custom theme colors.
        
        Args:
            primary_color: RGB tuple for primary color
            accent_color: RGB tuple for accent color
        """
        self.primary_color = RGBColor(*primary_color)
        self.accent_color = RGBColor(*accent_color)


def create_sample_devops_presentation() -> str:
    """
    Create a sample DevOps presentation to demonstrate the capabilities.
    
    Returns:
        Path to the created presentation file
    """
    ppt = PPTGenerator()
    
    # Set custom theme colors (DevOps blue theme)
    ppt.set_theme_colors((31, 73, 125), (79, 129, 189))
    
    # Title slide
    ppt.add_title_slide(
        "DevOps Best Practices 2025", 
        "Automation, CI/CD, and Infrastructure as Code",
        "DevOps Team"
    )
    
    # Introduction section
    ppt.add_section_slide("Introduction", (31, 73, 125))
    
    # Content slide - What is DevOps
    ppt.add_content_slide(
        "What is DevOps?",
        [
            "Cultural and professional movement focused on collaboration",
            "Emphasizes communication between development and operations",
            "Automates processes and infrastructure",
            "Improves deployment frequency and reliability",
            "Reduces time to market for new features"
        ]
    )
    
    # Two-column slide - Benefits vs Challenges
    ppt.add_two_column_slide(
        "DevOps: Benefits vs Challenges",
        [
            "Benefits:",
            "• Faster deployment cycles",
            "• Improved collaboration",
            "• Higher quality software",
            "• Better customer satisfaction",
            "• Reduced costs"
        ],
        [
            "Challenges:",
            "• Cultural resistance",
            "• Tool complexity",
            "• Security concerns",
            "• Skill gaps",
            "• Legacy system integration"
        ]
    )
    
    # CI/CD section
    ppt.add_section_slide("CI/CD Pipeline", (79, 129, 189))
    
    # Content slide - CI/CD stages
    ppt.add_content_slide(
        "CI/CD Pipeline Stages",
        [
            "Source Code Management (Git, GitLab, GitHub)",
            "Build Automation (Maven, Gradle, npm)",
            "Automated Testing (Unit, Integration, E2E)",
            "Code Quality Analysis (SonarQube, ESLint)",
            "Artifact Repository (Nexus, Artifactory)",
            "Deployment Automation (Jenkins, GitLab CI, GitHub Actions)",
            "Monitoring and Feedback (Prometheus, Grafana)"
        ]
    )
    
    # Chart slide - Deployment frequency
    chart_data = {
        "categories": ["Traditional", "Agile", "DevOps", "Elite DevOps"],
        "series": [
            {
                "name": "Deployments per Month",
                "values": [1, 4, 30, 200]
            }
        ]
    }
    ppt.add_chart_slide("Deployment Frequency Comparison", chart_data, "column")
    
    # Infrastructure section
    ppt.add_section_slide("Infrastructure as Code", (125, 73, 31))
    
    # Content slide - IaC tools
    ppt.add_content_slide(
        "Infrastructure as Code Tools",
        [
            "Terraform - Multi-cloud infrastructure provisioning",
            "AWS CloudFormation - AWS native IaC",
            "Azure Resource Manager - Azure native templates",
            "Ansible - Configuration management and automation",
            "Kubernetes - Container orchestration",
            "Docker - Containerization platform",
            "Helm - Kubernetes package manager"
        ]
    )
    
    # Monitoring section
    ppt.add_section_slide("Monitoring & Observability", (189, 129, 79))
    
    # Two-column slide - Monitoring tools
    ppt.add_two_column_slide(
        "Monitoring Stack",
        [
            "Metrics & Monitoring:",
            "• Prometheus",
            "• Grafana",
            "• DataDog",
            "• New Relic",
            "• CloudWatch"
        ],
        [
            "Logging & Tracing:",
            "• ELK Stack (Elasticsearch, Logstash, Kibana)",
            "• Fluentd",
            "• Jaeger",
            "• Zipkin",
            "• Splunk"
        ]
    )
    
    # Best practices
    ppt.add_content_slide(
        "DevOps Best Practices",
        [
            "Start small and iterate - Begin with pilot projects",
            "Automate everything - Reduce manual processes",
            "Implement comprehensive monitoring",
            "Foster collaboration culture",
            "Practice infrastructure as code",
            "Maintain security throughout (DevSecOps)",
            "Continuously learn and improve"
        ]
    )
    
    # Conclusion
    ppt.add_content_slide(
        "Key Takeaways",
        [
            "DevOps is a cultural transformation, not just tooling",
            "Start with small, manageable improvements",
            "Invest in automation and monitoring",
            "Focus on collaboration and communication",
            "Security should be integrated, not an afterthought",
            "Continuous learning and adaptation are essential"
        ]
    )
    
    # Thank you slide
    ppt.add_title_slide("Thank You!", "Questions & Discussion", "DevOps Team")
    
    # Save the presentation
    filename = f"devops_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    saved_path = ppt.save(filename)
    
    return saved_path


def main():
    """Main function to demonstrate PPT generation capabilities."""
    print("🎯 PowerPoint Presentation Generator")
    print("=" * 50)
    
    while True:
        print("\nChoose an option:")
        print("1. Create sample DevOps presentation")
        print("2. Create custom presentation")
        print("3. Exit")
        
        choice = input("\nEnter your choice (1-3): ").strip()
        
        if choice == "1":
            print("\n📊 Creating sample DevOps presentation...")
            try:
                saved_path = create_sample_devops_presentation()
                print(f"✅ Presentation created successfully!")
                print(f"📁 Saved as: {saved_path}")
                print(f"📄 File size: {os.path.getsize(saved_path) / 1024:.1f} KB")
            except Exception as e:
                print(f"❌ Error creating presentation: {e}")
        
        elif choice == "2":
            print("\n🎨 Creating custom presentation...")
            create_custom_presentation()
        
        elif choice == "3":
            print("\n👋 Goodbye!")
            break
        
        else:
            print("❌ Invalid choice. Please try again.")


def create_custom_presentation():
    """Interactive function to create a custom presentation."""
    ppt = PPTGenerator()
    
    # Get presentation details
    title = input("Enter presentation title: ").strip()
    subtitle = input("Enter subtitle (optional): ").strip()
    author = input("Enter author name (optional): ").strip()
    
    # Add title slide
    ppt.add_title_slide(title, subtitle, author)
    
    while True:
        print("\nAdd a slide:")
        print("1. Content slide with bullet points")
        print("2. Two-column slide")
        print("3. Section divider slide")
        print("4. Finish and save")
        
        slide_choice = input("Choose slide type (1-4): ").strip()
        
        if slide_choice == "1":
            slide_title = input("Enter slide title: ").strip()
            print("Enter bullet points (empty line to finish):")
            content = []
            while True:
                point = input("• ").strip()
                if not point:
                    break
                content.append(point)
            
            if content:
                ppt.add_content_slide(slide_title, content)
                print("✅ Content slide added!")
        
        elif slide_choice == "2":
            slide_title = input("Enter slide title: ").strip()
            
            print("Enter left column content (empty line to finish):")
            left_content = []
            while True:
                point = input("• ").strip()
                if not point:
                    break
                left_content.append(point)
            
            print("Enter right column content (empty line to finish):")
            right_content = []
            while True:
                point = input("• ").strip()
                if not point:
                    break
                right_content.append(point)
            
            if left_content or right_content:
                ppt.add_two_column_slide(slide_title, left_content, right_content)
                print("✅ Two-column slide added!")
        
        elif slide_choice == "3":
            section_title = input("Enter section title: ").strip()
            if section_title:
                ppt.add_section_slide(section_title)
                print("✅ Section slide added!")
        
        elif slide_choice == "4":
            filename = input("Enter filename (without .pptx): ").strip()
            if not filename:
                filename = f"custom_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            
            try:
                saved_path = ppt.save(filename)
                print(f"✅ Presentation saved successfully!")
                print(f"📁 Saved as: {saved_path}")
                print(f"📄 Slides: {ppt.get_slide_count()}")
                print(f"📄 File size: {os.path.getsize(saved_path) / 1024:.1f} KB")
                break
            except Exception as e:
                print(f"❌ Error saving presentation: {e}")
        
        else:
            print("❌ Invalid choice. Please try again.")


if __name__ == "__main__":
    main()